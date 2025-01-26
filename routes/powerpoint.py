# ============================================================
# FILE: routes/powerpoint.py
# ============================================================
import io
import os
from datetime import datetime

# python-pptx installieren: pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

# ------------------------------------------------------------
# Hilfsfunktionen für Textboxen, Tabellen usw.
# ------------------------------------------------------------

def _add_centered_textbox(slide, left, top, width, height, text, font_size=28, bold=False):
    """
    Fügt ein zentriertes Textfeld auf der Folie ein (Paragraph).
    """
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER

def _add_left_textbox(slide, left, top, width, height, text, font_size=14, bold=False):
    """
    Fügt ein links-ausgerichtetes Textfeld ein.
    """
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = PP_ALIGN.LEFT

# ------------------------------------------------------------
# HAUPTFUNKTION: export_baugruppe_pptx
# ------------------------------------------------------------

def export_baugruppe_pptx(tab1_data, tab2_data, tab3_steps, tab4_summary, filename="ParetoKalk_Presentation.pptx"):
    """
    Erstellt eine PPTX-Präsentation mit 3 Slides:
      1) Deckblatt: Titel, Logo, Projektname, Datum
      2) Kalkulation (Tabellen, Summen etc.)
      3) Vergleich (Unsere Kalkulation vs. Lieferant), inkl. Balkendiagramm

    Rückgabe: (BytesIO, filename) => per send_file downloadbar.
    """
    # In-Memory-Stream
    output = io.BytesIO()

    # 1) Neue Präsentation
    prs = Presentation()  # Standard 16:9 Folien

    # ---------------------------------------------------
    # SLIDE 1: Deckblatt
    # ---------------------------------------------------
    # Blank Layout (Index 6, i. d. R. bei Standard-Theme)
    slide1_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(slide1_layout)

    # Großer Titel
    _add_centered_textbox(
        slide1,
        left=Inches(1.0),
        top=Inches(1.0),
        width=Inches(8.0),
        height=Inches(1.0),
        text="Pareto-Kalk – Management One-Pager",
        font_size=34,
        bold=True
    )

    # Projektinfos
    project_name = tab1_data.get("projectName", "Unbekanntes Projekt")
    part_name = tab1_data.get("partName", "Bauteil X")
    date_str = datetime.now().strftime("%d.%m.%Y %H:%M")

    info_text = (
        f"Projekt: {project_name}\n"
        f"Bauteil: {part_name}\n"
        f"Erzeugt am: {date_str}\n\n"
        "Powered by python-pptx"
    )
    _add_left_textbox(
        slide1,
        left=Inches(1.0),
        top=Inches(2.2),
        width=Inches(6.0),
        height=Inches(2.0),
        text=info_text,
        font_size=16,
        bold=False
    )

    # Logo (optional)
    # Passe den Pfad an, falls dein Logo woanders liegt
    logo_path = os.path.join("static", "img", "logo.png")
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(
            logo_path,
            left=Inches(7.0),
            top=Inches(0.5),
            width=Inches(2.0),  # Skalierung, optional
            height=None
        )

    # ---------------------------------------------------
    # SLIDE 2: Kalkulations-Übersicht
    # ---------------------------------------------------
    slide2_layout = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(slide2_layout)

    _add_centered_textbox(
        slide2,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.6),
        "Kalkulations-Details (Tab1–4)",
        font_size=28,
        bold=True
    )

    # LINKS: Tabelle für Material & Overhead (Tab1/Tab2)
    rows = 7
    cols = 2
    table_left = Inches(0.5)
    table_top = Inches(1.3)
    table_width = Inches(4.5)
    table_height = Inches(3.0)

    shape_table = slide2.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = shape_table.table

    # Kopfzeile
    table.cell(0, 0).text = "Parameter"
    table.cell(0, 1).text = "Wert"

    # Zeile 1: Materialname
    mat_name = tab2_data.get("matName", "Alu")
    table.cell(1, 0).text = "Material"
    table.cell(1, 1).text = mat_name

    # Zeile 2: Preis
    mat_price = tab2_data.get("matPrice", 0)
    table.cell(2, 0).text = "Preis (€/kg)"
    table.cell(2, 1).text = f"{mat_price:.2f}"

    # Zeile 3: Gewicht
    mat_weight = tab2_data.get("matWeight", 0)
    table.cell(3, 0).text = "Gewicht (kg)"
    table.cell(3, 1).text = f"{mat_weight:.3f}"

    # Zeile 4: Fremdzukauf
    fremd_val = tab2_data.get("fremdValue", 0)
    table.cell(4, 0).text = "Fremdzukauf (€)"
    table.cell(4, 1).text = f"{fremd_val:.2f}"

    # Zeile 5: SG&A
    sga_pct = tab1_data.get("sgaPct", 0)
    table.cell(5, 0).text = "SG&A (%)"
    table.cell(5, 1).text = f"{sga_pct:.1f}%"

    # Zeile 6: Profit
    profit_pct = tab1_data.get("profitPct", 0)
    table.cell(6, 0).text = "Profit (%)"
    table.cell(6, 1).text = f"{profit_pct:.1f}%"

    # RECHTS: Fertigungstabelle (max 8 Steps) => tab3_steps
    fert_rows = len(tab3_steps[:8]) + 1  # +1 für Kopf
    fert_cols = 4
    fert_left = Inches(5.0)
    fert_top = Inches(1.3)
    fert_width = Inches(5.0)
    fert_height = Inches(3.0)

    fert_shape = slide2.shapes.add_table(fert_rows, fert_cols, fert_left, fert_top, fert_width, fert_height)
    fert_table = fert_shape.table

    fert_table.cell(0, 0).text = "Step"
    fert_table.cell(0, 1).text = "Zyklus (s)"
    fert_table.cell(0, 2).text = "Kosten/100"
    fert_table.cell(0, 3).text = "CO₂/100"

    for i, step in enumerate(tab3_steps[:8], start=1):
        name = step.get("stepName", f"Step{i}")
        cyc = step.get("cycTime", 0)
        c100 = step.get("kosten_100", 0)
        co2 = step.get("co2_100", 0)

        fert_table.cell(i, 0).text = name
        fert_table.cell(i, 1).text = f"{cyc:.1f}"
        fert_table.cell(i, 2).text = f"{c100:.2f}"
        fert_table.cell(i, 3).text = f"{co2:.2f}"

    # UNTEN: Summary (Tab4)
    summary_text = ""
    summary_text += f"Mat-Einzelkosten/100: {tab4_summary.get('matEinzel', 0):.2f} €\n"
    summary_text += f"Fremdzukauf/100: {tab4_summary.get('fremd', 0):.2f} €\n"
    summary_text += f"Herstellkosten/100: {tab4_summary.get('herstell', 0):.2f} €\n"
    summary_text += f"SG&A: {tab4_summary.get('sga', 0):.2f} € | Profit: {tab4_summary.get('profit', 0):.2f} €\n"
    summary_text += f"Total/100: {tab4_summary.get('total', 0):.2f} €\n"
    summary_text += f"CO₂/100: {tab4_summary.get('co2_100', 0):.2f} kg"

    _add_left_textbox(
        slide2,
        left=Inches(0.5),
        top=Inches(4.6),
        width=Inches(9.0),
        height=Inches(2.0),
        text="**Kalkulations-Ergebnis**\n" + summary_text,
        font_size=14,
        bold=False
    )

    # ---------------------------------------------------
    # SLIDE 3: Vergleich (Unsere Kalkulation vs. Lieferant)
    # ---------------------------------------------------
    slide3_layout = prs.slide_layouts[6]
    slide3 = prs.slides.add_slide(slide3_layout)

    _add_centered_textbox(
        slide3,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.6),
        "Vergleich: Unsere Kalkulation vs. Lieferant",
        font_size=24,
        bold=True
    )

    # a) Wir nehmen an, du hast fiktive Lieferantenwerte
    #    Du kannst sie in tab4_summary["supplier"] ablegen, oder fix
    #    Hier demonstrativ mit Mock-Daten:
    supplier_data = {
        "total": 280.0,  # z.B. vs. unser total=250
        "co2_100": 15.0  # vs. unser=12
    }
    # b) Erzeuge ein Balkendiagramm: 2 Kategorien (Kosten, CO2), 2 Serien (Wir, Lieferant)

    chart_data = CategoryChartData()
    chart_data.categories = ["Kosten/100 (€)", "CO2/100 (kg)"]

    # Serie1: Unsere Kalkulation
    our_total = tab4_summary.get("total", 0)
    our_co2 = tab4_summary.get("co2_100", 0)
    chart_data.add_series("Unsere Kalkulation", (our_total, our_co2))

    # Serie2: Lieferant
    sup_total = supplier_data["total"]
    sup_co2 = supplier_data["co2_100"]
    chart_data.add_series("Lieferant", (sup_total, sup_co2))

    x = Inches(0.5)
    y = Inches(1.5)
    cx = Inches(5.0)
    cy = Inches(4.0)
    graphic_frame = slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x, y, cx, cy,
        chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.value_axis.has_major_gridlines = False

    # c) Tabelle: Detaillierter Vergleich
    comp_left = Inches(6.0)
    comp_top = Inches(1.5)
    comp_width = Inches(4.0)
    comp_height = Inches(2.5)
    comp_table_shape = slide3.shapes.add_table(4, 3, comp_left, comp_top, comp_width, comp_height)
    comp_table = comp_table_shape.table

    # Kopf:
    comp_table.cell(0, 0).text = "Position"
    comp_table.cell(0, 1).text = "Wir"
    comp_table.cell(0, 2).text = "Lieferant"

    # Zeile 1:
    comp_table.cell(1, 0).text = "Total/100 (€)"
    comp_table.cell(1, 1).text = f"{our_total:.2f}"
    comp_table.cell(1, 2).text = f"{sup_total:.2f}"

    # Zeile 2:
    comp_table.cell(2, 0).text = "CO₂/100 (kg)"
    comp_table.cell(2, 1).text = f"{our_co2:.2f}"
    comp_table.cell(2, 2).text = f"{sup_co2:.2f}"

    # Zeile 3:
    # Optional Delta
    delta_cost = sup_total - our_total
    delta_co2 = sup_co2 - our_co2
    comp_table.cell(3, 0).text = "Delta"
    comp_table.cell(3, 1).text = f"-" if delta_cost < 0 else f"+{delta_cost:.2f}€"
    comp_table.cell(3, 2).text = f"+{delta_co2:.2f} kg" if delta_co2 >= 0 else f"{delta_co2:.2f} kg"

    # pptx speichern
    prs.save(output)
    output.seek(0)
    return output, filename

# ------------------------------------------------------------
# ENDE FILE: routes/powerpoint.py
# ------------------------------------------------------------
